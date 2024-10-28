import os
import sys
from pptx import Presentation
from pptx.util import Inches

def add_images_from_folder(prs, folder_path):
    """Add PNG images from a folder as slides to the PowerPoint presentation."""
    png_files = sorted([f for f in os.listdir(folder_path) if f.endswith(".png")])

    for png_file in png_files:
        png_path = os.path.join(folder_path, png_file)

        # Add a blank slide layout (could be adjusted if needed)
        slide_layout = prs.slide_layouts[5]  # 5 is for a blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Add the image to the slide (resize it to fit the slide)
        left = top = Inches(0)  # Position image at the top-left corner
        slide.shapes.add_picture(png_path, left, top, width=prs.slide_width, height=prs.slide_height)

def create_presentation_from_folders(output_ppt, folders):
    """Create a PowerPoint presentation from a list of folders containing PNG slides."""
    # Create a new PowerPoint presentation
    prs = Presentation()

    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.33)  # Widescreen width for 16:9 (13.33 inches)
    prs.slide_height = Inches(7.5)   # Widescreen height for 16:9 (7.5 inches)

    # Iterate over each folder and add images from each
    for folder in folders:
        if not os.path.isdir(folder):
            print(f"Skipping {folder}, not a directory.")
            continue

        print(f"Adding images from folder: {folder}")
        add_images_from_folder(prs, folder)

    # Save the final PowerPoint presentation
    prs.save(output_ppt)
    print(f"PowerPoint file created: {output_ppt}")

if len(sys.argv) < 3:
    print("Usage: python create_ppt_from_folders.py <output_pptx> <folder1> <folder2> ...")
    sys.exit(1)

def main():
    output_pptx = sys.argv[1]
    folders = sys.argv[2:]

    # Create the PowerPoint presentation from the specified folders
    create_presentation_from_folders(output_pptx, folders)

if __name__ == "__main__":
    main()
