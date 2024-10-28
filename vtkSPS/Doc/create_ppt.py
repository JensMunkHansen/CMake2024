import os
import sys
from pptx import Presentation
from pptx.util import Inches

# Check if the correct number of arguments is provided
if len(sys.argv) != 3:
    print("Usage: python create_ppt.py <input_image_directory> <output_ppt_file>")
    sys.exit(1)

# Get input directory and output PowerPoint file from arguments
input_dir = sys.argv[1]
output_ppt = sys.argv[2]

# Create a new PowerPoint presentation
prs = Presentation()

# Set slide size to widescreen (16:9)
prs.slide_width = Inches(13.33)  # Widescreen width for 16:9 (13.33 inches)
prs.slide_height = Inches(7.5)   # Widescreen height for 16:9 (7.5 inches)

# Get a list of PNG files in the input directory, sorted alphabetically
png_files = sorted([f for f in os.listdir(input_dir) if f.endswith(".png")])

# Loop through the PNG files and add each as a new slide
for png_file in png_files:
    png_path = os.path.join(input_dir, png_file)
    
    # Add a blank slide layout (could be adjusted if needed)
    slide_layout = prs.slide_layouts[5]  # 5 is for a blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Add the image to the slide (full slide size assumed here)
    left = top = Inches(0)  # Position image at the top-left corner
    slide.shapes.add_picture(png_path, left, top, width=prs.slide_width, height=prs.slide_height)

# Save the PowerPoint presentation
prs.save(output_ppt)

print(f"PowerPoint file created: {output_ppt}")
